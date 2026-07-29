"""
PDF内の日本語だけを指定言語へ翻訳し、編集可能なテキストボックスを重ねたPPTXを生成する。

処理:
  1. Azure Document Intelligence (prebuilt-read) で日本語行と座標を取得
  2. Azure OpenAI で行ごとに指定言語へ翻訳
  3. PDFページを背景画像として配置
  4. 日本語行を背景色付きの編集可能な翻訳テキストボックスで覆う

CLI:
  python pdf_translate_to_pptx.py --input input.pdf --output output.pptx \
    --target-language en

stdout:
  JSON {"pages": N, "translatedLines": N, "engine": "..."}
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
from collections import Counter
from dataclasses import dataclass
from io import BytesIO
from typing import Any, Iterable

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.core.credentials import AzureKeyCredential
except ImportError as exc:
    raise RuntimeError(
        "azure-ai-documentintelligence is required for PDF translation"
    ) from exc


MAX_PAGES = 10
RENDER_DPI = 144
TRANSLATION_BATCH_SIZE = 60
JAPANESE_RE = re.compile(
    r"[\u3040-\u30ff\u31f0-\u31ff\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]"
)


@dataclass(frozen=True)
class TargetLanguage:
    code: str
    english_name: str
    japanese_name: str
    font_name: str


TARGET_LANGUAGES: dict[str, TargetLanguage] = {
    "en": TargetLanguage("en", "English", "英語", "Arial"),
    "pt": TargetLanguage("pt", "Portuguese", "ポルトガル語", "Arial"),
    "vi": TargetLanguage("vi", "Vietnamese", "ベトナム語", "Arial"),
    "id": TargetLanguage("id", "Indonesian", "インドネシア語", "Arial"),
    "zh-CN": TargetLanguage(
        "zh-CN", "Simplified Chinese", "中国語（簡体字）", "Microsoft YaHei"
    ),
    "ko": TargetLanguage("ko", "Korean", "韓国語", "Malgun Gothic"),
    "es": TargetLanguage("es", "Spanish", "スペイン語", "Arial"),
    "fil": TargetLanguage("fil", "Filipino (Tagalog)", "タガログ語", "Arial"),
}


@dataclass
class TextRegion:
    id: str
    page_index: int
    source: str
    polygon: list[tuple[float, float]]
    page_width: float
    page_height: float
    page_unit: str
    translation: str = ""

    @property
    def bbox(self) -> tuple[float, float, float, float]:
        xs = [point[0] for point in self.polygon]
        ys = [point[1] for point in self.polygon]
        return min(xs), min(ys), max(xs), max(ys)


def _env(name: str, fallback: str | None = None) -> str:
    value = os.environ.get(name, "").strip()
    if value:
        return value
    if fallback:
        return os.environ.get(fallback, "").strip()
    return ""


def _doc_intelligence_client() -> DocumentIntelligenceClient:
    endpoint = _env("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
    key = _env("AZURE_DOCUMENT_INTELLIGENCE_KEY")
    if not endpoint or not key:
        raise RuntimeError(
            "AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT/KEY are not configured"
        )
    return DocumentIntelligenceClient(
        endpoint=endpoint,
        credential=AzureKeyCredential(key),
    )


def _point_xy(point: Any) -> tuple[float, float]:
    if hasattr(point, "x") and hasattr(point, "y"):
        return float(point.x), float(point.y)
    if isinstance(point, dict):
        return float(point["x"]), float(point["y"])
    if isinstance(point, (list, tuple)) and len(point) >= 2:
        return float(point[0]), float(point[1])
    raise ValueError(f"Unsupported polygon point: {point!r}")


def _polygon_points(polygon: Any) -> list[tuple[float, float]]:
    raw = list(polygon or [])
    if raw and all(isinstance(item, (int, float)) for item in raw):
        if len(raw) % 2 != 0:
            raise ValueError(f"Polygon has an odd number of coordinates: {raw!r}")
        return [
            (float(raw[index]), float(raw[index + 1]))
            for index in range(0, len(raw), 2)
        ]
    return [_point_xy(point) for point in raw]


def _extract_regions(pdf_path: str) -> tuple[list[TextRegion], int]:
    client = _doc_intelligence_client()
    with open(pdf_path, "rb") as source:
        poller = client.begin_analyze_document(
            "prebuilt-read",
            body=source,
            content_type="application/pdf",
            locale="ja-JP",
        )
    result = poller.result()
    pages = list(result.pages or [])
    if not pages:
        raise RuntimeError("Document Intelligence returned no pages")
    if len(pages) > MAX_PAGES:
        raise RuntimeError(
            f"PDF has {len(pages)} pages; the current limit is {MAX_PAGES} pages"
        )

    regions: list[TextRegion] = []
    for page_index, page in enumerate(pages):
        width = float(page.width or 0)
        height = float(page.height or 0)
        unit = str(page.unit or "inch")
        if width <= 0 or height <= 0:
            continue
        for line_index, line in enumerate(page.lines or []):
            source_text = (line.content or "").strip()
            polygon = _polygon_points(line.polygon)
            if not source_text or not JAPANESE_RE.search(source_text):
                continue
            if len(polygon) < 4:
                continue
            regions.append(
                TextRegion(
                    id=f"p{page_index + 1}-l{line_index + 1}",
                    page_index=page_index,
                    source=source_text,
                    polygon=polygon,
                    page_width=width,
                    page_height=height,
                    page_unit=unit,
                )
            )
    return regions, len(pages)


def _openai_settings() -> tuple[str, str, str, str]:
    key = _env("AZURE_OPENAI_API_KEY", "AZURE_OPENAI_VISION_API_KEY")
    instance = _env(
        "AZURE_OPENAI_API_INSTANCE_NAME",
        "AZURE_OPENAI_VISION_API_INSTANCE_NAME",
    )
    deployment = _env(
        "AZURE_OPENAI_API_DEPLOYMENT_NAME",
        "AZURE_OPENAI_VISION_API_DEPLOYMENT_NAME",
    )
    api_version = _env(
        "AZURE_OPENAI_API_VERSION",
        "AZURE_OPENAI_VISION_API_VERSION",
    ) or "2024-12-01-preview"
    if not key or not instance or not deployment:
        raise RuntimeError("Azure OpenAI environment variables are not configured")
    return key, instance, deployment, api_version


def _extract_json_object(value: str) -> dict[str, Any]:
    cleaned = value.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    try:
        parsed = json.loads(cleaned)
        if isinstance(parsed, dict):
            return parsed
    except json.JSONDecodeError:
        pass
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start >= 0 and end > start:
        parsed = json.loads(cleaned[start : end + 1])
        if isinstance(parsed, dict):
            return parsed
    raise ValueError("Azure OpenAI response did not contain a JSON object")


def _translate_batch(
    regions: list[TextRegion],
    target_language: TargetLanguage,
) -> dict[str, str]:
    key, instance, deployment, api_version = _openai_settings()
    endpoint = (
        f"https://{instance}.openai.azure.com/openai/deployments/"
        f"{urllib.parse.quote(deployment, safe='')}/chat/completions"
        f"?api-version={urllib.parse.quote(api_version, safe='')}"
    )
    items = [{"id": region.id, "text": region.source} for region in regions]
    request_body = {
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are a professional translator for public information brochures. "
                    f"Translate every item's Japanese text into {target_language.english_name}. "
                    f"Output all translated wording in {target_language.english_name} only. "
                    "Preserve numbers, dates, phone numbers, URLs, list markers, and legal "
                    "meaning. Use natural, concise wording that fits in the original layout. "
                    "If a line is a fragment because the source wraps, translate it as a "
                    "matching fragment using nearby items as context. Return JSON only in "
                    "this exact shape: "
                    '{"translations":[{"id":"...","text":"..."}]}. '
                    "Return exactly one result for every input id."
                ),
            },
            {
                "role": "user",
                "content": json.dumps({"items": items}, ensure_ascii=False),
            },
        ],
        "temperature": 0.1,
        "response_format": {"type": "json_object"},
    }
    request = urllib.request.Request(
        endpoint,
        data=json.dumps(request_body, ensure_ascii=False).encode("utf-8"),
        headers={"Content-Type": "application/json", "api-key": key},
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=180) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(
            f"Azure OpenAI translation failed: HTTP {exc.code}: {body[:500]}"
        ) from exc

    content = payload["choices"][0]["message"]["content"]
    parsed = _extract_json_object(content)
    translations = parsed.get("translations")
    if not isinstance(translations, list):
        raise RuntimeError("Azure OpenAI response is missing translations")
    output: dict[str, str] = {}
    for item in translations:
        if not isinstance(item, dict):
            continue
        item_id = str(item.get("id", "")).strip()
        text = str(item.get("text", "")).strip()
        if item_id and text:
            output[item_id] = text
    return output


def _translate_regions(
    regions: list[TextRegion],
    target_language: TargetLanguage,
) -> None:
    for start in range(0, len(regions), TRANSLATION_BATCH_SIZE):
        batch = regions[start : start + TRANSLATION_BATCH_SIZE]
        translated = _translate_batch(batch, target_language)
        missing: list[TextRegion] = []
        for region in batch:
            region.translation = translated.get(region.id, "").strip()
            if not region.translation:
                missing.append(region)
        if missing:
            print(
                "[pdf-translate] retrying omitted translations: "
                + ", ".join(region.id for region in missing[:10]),
                file=sys.stderr,
            )
            retry = _translate_batch(missing, target_language)
            for region in missing:
                region.translation = retry.get(region.id, "").strip()
            still_missing = [
                region.id for region in missing if not region.translation
            ]
            if still_missing:
                raise RuntimeError(
                    "Azure OpenAI omitted translations after retry: "
                    + ", ".join(still_missing[:10])
                )


def _render_pages(pdf_path: str) -> list[Image.Image]:
    document = fitz.open(pdf_path)
    if len(document) > MAX_PAGES:
        document.close()
        raise RuntimeError(
            f"PDF has {len(document)} pages; the current limit is {MAX_PAGES} pages"
        )
    scale = RENDER_DPI / 72.0
    images: list[Image.Image] = []
    try:
        for page in document:
            pixmap = page.get_pixmap(
                matrix=fitz.Matrix(scale, scale),
                alpha=False,
            )
            image = Image.frombytes(
                "RGB",
                (pixmap.width, pixmap.height),
                pixmap.samples,
            )
            images.append(image)
    finally:
        document.close()
    return images


def _quantize(rgb: tuple[int, int, int], step: int = 16) -> tuple[int, int, int]:
    return tuple(min(255, int(round(channel / step) * step)) for channel in rgb)


def _rgb_distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return math.sqrt(sum((a[index] - b[index]) ** 2 for index in range(3)))


def _sample_colors(
    image: Image.Image,
    bbox_px: tuple[int, int, int, int],
) -> tuple[tuple[int, int, int], tuple[int, int, int]]:
    x0, y0, x1, y1 = bbox_px
    width, height = image.size
    x0 = max(0, min(width - 1, x0))
    y0 = max(0, min(height - 1, y0))
    x1 = max(x0 + 1, min(width, x1))
    y1 = max(y0 + 1, min(height, y1))
    pad = max(2, min(8, int(max(x1 - x0, y1 - y0) * 0.08)))

    border_pixels: list[tuple[int, int, int]] = []
    for y in range(max(0, y0 - pad), min(height, y1 + pad)):
        for x in range(max(0, x0 - pad), min(width, x1 + pad)):
            if x0 <= x < x1 and y0 <= y < y1:
                continue
            border_pixels.append(image.getpixel((x, y)))
    if not border_pixels:
        border_pixels = [(255, 255, 255)]
    bg_counter = Counter(_quantize(pixel) for pixel in border_pixels)
    background = bg_counter.most_common(1)[0][0]

    inside = image.crop((x0, y0, x1, y1))
    stride = max(1, int(math.sqrt(max(1, inside.width * inside.height) / 5000)))
    foreground_candidates: Counter[tuple[int, int, int]] = Counter()
    pixels = inside.load()
    for y in range(0, inside.height, stride):
        for x in range(0, inside.width, stride):
            color = _quantize(pixels[x, y])
            if _rgb_distance(color, background) >= 70:
                foreground_candidates[color] += 1

    if foreground_candidates:
        foreground = max(
            foreground_candidates,
            key=lambda color: (
                foreground_candidates[color] * min(220, _rgb_distance(color, background)),
                _rgb_distance(color, background),
            ),
        )
    else:
        luminance = (
            0.2126 * background[0]
            + 0.7152 * background[1]
            + 0.0722 * background[2]
        )
        foreground = (0, 0, 0) if luminance > 145 else (255, 255, 255)
    return background, foreground


def _region_to_pixels(
    region: TextRegion,
    image: Image.Image,
) -> tuple[int, int, int, int]:
    x0, y0, x1, y1 = region.bbox
    scale_x = image.width / region.page_width
    scale_y = image.height / region.page_height
    return (
        int(x0 * scale_x),
        int(y0 * scale_y),
        int(math.ceil(x1 * scale_x)),
        int(math.ceil(y1 * scale_y)),
    )


def _image_bytes(image: Image.Image) -> BytesIO:
    stream = BytesIO()
    image.save(stream, format="PNG", optimize=True)
    stream.seek(0)
    return stream


def _font_size_points(
    width_in: float,
    height_in: float,
    text: str,
    rotated: bool,
) -> float:
    logical_width = height_in if rotated else width_in
    logical_height = width_in if rotated else height_in
    text_width = max(
        1.0,
        sum(
            0.0
            if unicodedata.combining(character)
            else 2.0
            if unicodedata.east_asian_width(character) in {"W", "F"}
            else 1.0
            for character in text
        ),
    )
    lines = max(1, math.ceil(text_width / max(4, int(logical_width * 10))))
    by_height = logical_height * 72 * 0.78 / lines
    by_width = logical_width * 72 * 1.85 / max(4.0, text_width)
    return max(7.0, min(30.0, by_height, by_width))


def _set_run_typeface(run: Any, font_name: str) -> None:
    run.font.name = font_name
    properties = run._r.get_or_add_rPr()
    for tag_name in ("a:latin", "a:ea", "a:cs"):
        element = properties.find(qn(tag_name))
        if element is None:
            element = OxmlElement(tag_name)
            properties.append(element)
        element.set("typeface", font_name)


def _build_presentation(
    pdf_path: str,
    output_path: str,
    regions: list[TextRegion],
    target_language: TargetLanguage,
) -> None:
    pages = _render_pages(pdf_path)
    if not pages:
        raise RuntimeError("PDF contains no pages")

    with fitz.open(pdf_path) as document:
        first_rect = document[0].rect
        slide_width_in = first_rect.width / 72.0
        slide_height_in = first_rect.height / 72.0

    presentation = Presentation()
    presentation.slide_width = Inches(slide_width_in)
    presentation.slide_height = Inches(slide_height_in)
    blank_layout = presentation.slide_layouts[6]
    regions_by_page: dict[int, list[TextRegion]] = {}
    for region in regions:
        regions_by_page.setdefault(region.page_index, []).append(region)

    for page_index, image in enumerate(pages):
        slide = presentation.slides.add_slide(blank_layout)
        page_regions: list[
            tuple[
                TextRegion,
                tuple[float, float, float, float],
                tuple[int, int, int],
                tuple[int, int, int],
            ]
        ] = []
        cleaned_image = image.copy()
        draw = ImageDraw.Draw(cleaned_image)

        for region in regions_by_page.get(page_index, []):
            x0, y0, x1, y1 = region.bbox
            scale_x = slide_width_in / region.page_width
            scale_y = slide_height_in / region.page_height
            left = max(0.0, x0 * scale_x)
            top = max(0.0, y0 * scale_y)
            width = max(0.06, (x1 - x0) * scale_x)
            height = max(0.06, (y1 - y0) * scale_y)

            pad = min(0.025, max(0.008, height * 0.08))
            left = max(0.0, left - pad)
            top = max(0.0, top - pad)
            width = min(slide_width_in - left, width + pad * 2)
            height = min(slide_height_in - top, height + pad * 2)
            background, foreground = _sample_colors(
                image,
                _region_to_pixels(region, image),
            )
            erase_box = (
                max(0, int(left / slide_width_in * cleaned_image.width)),
                max(0, int(top / slide_height_in * cleaned_image.height)),
                min(
                    cleaned_image.width,
                    int(math.ceil((left + width) / slide_width_in * cleaned_image.width)),
                ),
                min(
                    cleaned_image.height,
                    int(math.ceil((top + height) / slide_height_in * cleaned_image.height)),
                ),
            )
            draw.rectangle(erase_box, fill=background)
            page_regions.append(
                (region, (left, top, width, height), background, foreground)
            )

        slide.shapes.add_picture(
            _image_bytes(cleaned_image),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

        for region, box, _background, foreground in page_regions:
            left, top, width, height = box
            shape = slide.shapes.add_textbox(
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
            shape.name = f"{target_language.english_name} translation {region.id}"
            shape.fill.background()
            shape.line.fill.background()

            rotated = height > width * 1.2
            if rotated:
                # Keep the original narrow vertical box and lay the translation sideways.
                shape.text_frame._txBody.bodyPr.set("vert", "vert270")

            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.015)
            text_frame.margin_right = Inches(0.015)
            text_frame.margin_top = Inches(0.005)
            text_frame.margin_bottom = Inches(0.005)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            run = paragraph.add_run()
            run.text = region.translation
            _set_run_typeface(run, target_language.font_name)
            run.font.size = Pt(
                _font_size_points(width, height, region.translation, rotated)
            )
            run.font.bold = height >= 0.22 or rotated
            run.font.color.rgb = RGBColor(*foreground)

    presentation.save(output_path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument(
        "--target-language",
        choices=TARGET_LANGUAGES.keys(),
        default="en",
        help="Translation target language code (default: en)",
    )
    args = parser.parse_args()
    target_language = TARGET_LANGUAGES[args.target_language]

    regions, page_count = _extract_regions(args.input)
    if not regions:
        raise RuntimeError("No Japanese text was detected in the PDF")
    print(
        f"[pdf-translate] detected {len(regions)} Japanese lines on {page_count} pages",
        file=sys.stderr,
    )
    _translate_regions(regions, target_language)
    _build_presentation(args.input, args.output, regions, target_language)
    print(
        json.dumps(
            {
                "pages": page_count,
                "translatedLines": len(regions),
                "targetLanguage": target_language.code,
                "targetLanguageName": target_language.japanese_name,
                "engine": "azure_document_intelligence+azure_openai+python_pptx",
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:
        print(f"[pdf-translate] ERROR: {exc}", file=sys.stderr)
        raise
