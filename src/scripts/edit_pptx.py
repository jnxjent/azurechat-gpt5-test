import argparse
import json
import math
import sys
from copy import deepcopy
from pathlib import Path
from typing import Any
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Pt
from lxml import etree


def normalize_hex(value: str | None) -> str | None:
    if not value:
        return None
    normalized = str(value).replace("#", "").strip().upper()
    if len(normalized) == 6 and all(ch in "0123456789ABCDEF" for ch in normalized):
        return normalized
    return None


def rgb_to_hsl(rgb: tuple[int, int, int]) -> tuple[float, float, float]:
    r, g, b = [c / 255.0 for c in rgb]
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    l = (max_c + min_c) / 2
    if max_c == min_c:
        return 0.0, 0.0, l
    d = max_c - min_c
    s = d / (1 - abs(2 * l - 1))
    if max_c == r:
        h = 60 * (((g - b) / d) % 6)
    elif max_c == g:
        h = 60 * (((b - r) / d) + 2)
    else:
        h = 60 * (((r - g) / d) + 4)
    return h % 360, s, l


def hsl_to_rgb(h: float, s: float, l: float) -> tuple[int, int, int]:
    c = (1 - abs(2 * l - 1)) * s
    hh = h / 60
    x = c * (1 - abs((hh % 2) - 1))
    if 0 <= hh < 1:
        r1, g1, b1 = c, x, 0
    elif hh < 2:
        r1, g1, b1 = x, c, 0
    elif hh < 3:
        r1, g1, b1 = 0, c, x
    elif hh < 4:
        r1, g1, b1 = 0, x, c
    elif hh < 5:
        r1, g1, b1 = x, 0, c
    else:
        r1, g1, b1 = c, 0, x
    m = l - c / 2
    return (
        round((r1 + m) * 255),
        round((g1 + m) * 255),
        round((b1 + m) * 255),
    )


def recolor_preserving_tone(rgb: tuple[int, int, int], target_hex: str) -> str:
    hue, sat, lig = rgb_to_hsl(rgb)
    target_rgb = tuple(int(target_hex[i : i + 2], 16) for i in range(0, 6, 2))
    target_hue, _, _ = rgb_to_hsl(target_rgb)
    new_rgb = hsl_to_rgb(target_hue, max(sat, 0.35), lig)
    return "".join(f"{c:02X}" for c in new_rgb)


def shift_lightness(target_hex: str, lightness_delta: float) -> str:
    target_rgb = tuple(int(target_hex[i : i + 2], 16) for i in range(0, 6, 2))
    hue, sat, lig = rgb_to_hsl(target_rgb)
    next_lig = max(0.08, min(0.92, lig + lightness_delta))
    adjusted = hsl_to_rgb(hue, max(sat, 0.35), next_lig)
    return "".join(f"{c:02X}" for c in adjusted)


def is_neutral(rgb: tuple[int, int, int]) -> bool:
    hue, sat, lig = rgb_to_hsl(rgb)
    return sat < 0.12 or lig < 0.08 or lig > 0.95


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def apply_font_face(shape, font_face: str) -> bool:
    changed = False
    if not getattr(shape, "has_text_frame", False):
        return False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name != font_face:
                run.font.name = font_face
                changed = True
    return changed


def replace_text(shape, replacements: list[dict]) -> bool:
    """テキストを置換する。shape.text = ... は書式・枠をリセットするため使わず、
    各 run のテキストを直接置換して書式を保持する。
    appendToRun が指定された場合は find のテキストを持つ run の末尾にテキストを追記する。"""
    if not getattr(shape, "has_text_frame", False):
        return False

    changed = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_run = run.text
            updated_run = original_run
            for item in replacements:
                find = item.get("find") or ""
                append_text = item.get("appendToRun") or ""
                rep = item.get("replace") or ""
                if find:
                    if append_text:
                        # appendToRun: find で run を特定し、末尾にテキストを追記
                        if find in updated_run:
                            updated_run = updated_run.replace(find, find + append_text, 1)
                    elif rep:
                        updated_run = updated_run.replace(find, rep)
            if updated_run != original_run:
                run.text = updated_run
                changed = True
    return changed


def try_get_rgb(color_format) -> tuple[int, int, int] | None:
    try:
        if color_format.type == MSO_COLOR_TYPE.RGB and color_format.rgb is not None:
            rgb = str(color_format.rgb)
            return tuple(int(rgb[i : i + 2], 16) for i in range(0, 6, 2))
    except Exception:
        return None
    return None


def _palette_hex_for_rgb(rgb: tuple[int, int, int], palette: dict) -> str | None:
    """明度でセマンティックロールを判定してパレット色を返す。
    is_neutral() で lig>0.95 (ほぼ白) と sat<0.12 は呼び元でフィルタ済み。"""
    _, _, lig = rgb_to_hsl(rgb)
    if lig > 0.75:
        # 淡色背景 (sectionBg 相当) — 完全な白は is_neutral() で除去済み
        return palette.get("sectionBg")
    return palette.get("accentA" if lig < 0.45 else "accentB")


def _text_hex_for_rgb(rgb: tuple[int, int, int], palette: dict) -> str | None:
    """テキストランの色をパレットにマッピングする。
    白文字・低彩度グレーは呼び元の is_neutral() でフィルタ済み。"""
    _, _, lig = rgb_to_hsl(rgb)
    if lig < 0.45:
        return palette.get("bodyText")   # 暗色テキスト → bodyText
    return palette.get("accentB")        # 中間明度の強調テキスト → accentB


def recolor_fill(shape, target_hex: str | None, palette: dict | None = None) -> bool:
    try:
        fill = shape.fill
        if fill.type != MSO_FILL_TYPE.SOLID:
            return False
        rgb = try_get_rgb(fill.fore_color)
        if not rgb or is_neutral(rgb):
            return False
        if palette:
            new_hex = _palette_hex_for_rgb(rgb, palette)
            if not new_hex:
                return False
        elif target_hex:
            new_hex = recolor_preserving_tone(rgb, target_hex)
        else:
            return False
        fill.fore_color.rgb = RGBColor.from_string(new_hex)
        return True
    except Exception:
        return False


def recolor_line(shape, target_hex: str | None, palette: dict | None = None) -> bool:
    """既存の枠線色を直接 XML で書き換える。
    python-pptx の shape.line アクセサは使わない
    （アクセスだけで <a:ln> が生成される／<a:noFill> が <a:solidFill> に変換されるため）。
    <a:ln><a:solidFill><a:srgbClr> が明示されている場合のみ処理。"""
    try:
        _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        el = getattr(shape, "element", None)
        if el is None:
            return False
        ln_el = el.find(f".//{{{_ns}}}ln")
        if ln_el is None:
            return False
        solid = ln_el.find(f"{{{_ns}}}solidFill")
        if solid is None:
            return False  # noFill や空の <a:ln> はスキップ
        srgb = solid.find(f"{{{_ns}}}srgbClr")
        if srgb is None:
            return False  # schemeClr 等はスキップ
        val = srgb.get("val", "")
        if len(val) != 6:
            return False
        try:
            rgb = tuple(int(val[i : i + 2], 16) for i in range(0, 6, 2))
        except ValueError:
            return False
        if is_neutral(rgb):
            return False
        if palette:
            new_hex = _palette_hex_for_rgb(rgb, palette)
            if not new_hex:
                return False
        elif target_hex:
            new_hex = recolor_preserving_tone(rgb, target_hex)
        else:
            return False
        srgb.set("val", new_hex)
        return True
    except Exception:
        return False


def recolor_slide_background(slide, target_hex: str | None, palette: dict | None = None) -> bool:
    """Recolor the slide's background fill (<p:bg>) if it is a non-neutral solid color."""
    try:
        background = slide.background
        fill = background.fill
        if fill.type != MSO_FILL_TYPE.SOLID:
            return False
        rgb = try_get_rgb(fill.fore_color)
        if not rgb or is_neutral(rgb):
            return False
        if palette:
            new_hex = _palette_hex_for_rgb(rgb, palette)
            if not new_hex:
                return False
        elif target_hex:
            new_hex = recolor_preserving_tone(rgb, target_hex)
        else:
            return False
        fill.fore_color.rgb = RGBColor.from_string(new_hex)
        return True
    except Exception:
        return False


def recolor_table(shape, target_hex: str | None, palette: dict | None = None) -> bool:
    if not hasattr(shape, "has_table") or not shape.has_table:
        return False
    changed = False
    for row in shape.table.rows:
        for cell in row.cells:
            try:
                fill = cell.fill
                if fill.type == MSO_FILL_TYPE.SOLID:
                    rgb = try_get_rgb(fill.fore_color)
                    if rgb and not is_neutral(rgb):
                        if palette:
                            new_hex = _palette_hex_for_rgb(rgb, palette)
                            if not new_hex:
                                continue
                        elif target_hex:
                            new_hex = recolor_preserving_tone(rgb, target_hex)
                        else:
                            continue
                        fill.fore_color.rgb = RGBColor.from_string(new_hex)
                        changed = True
            except Exception:
                pass
    return changed


def recolor_text_runs(shape, target_hex: str | None, palette: dict | None = None) -> int:
    """shape のテキストランの明示的 RGB 色を新パレット / アクセント色に置換。
    テーブルセル内テキストも対象。戻り値: 変更したランの数。"""

    def _recolor_frame(text_frame) -> int:
        cnt = 0
        for para in text_frame.paragraphs:
            for run in para.runs:
                try:
                    fc = run.font.color
                    if fc.type != MSO_COLOR_TYPE.RGB:
                        continue
                    rgb = try_get_rgb(fc)
                    if not rgb or is_neutral(rgb):
                        continue
                    if palette:
                        new_hex = _text_hex_for_rgb(rgb, palette)
                    elif target_hex:
                        new_hex = recolor_preserving_tone(rgb, target_hex)
                    else:
                        continue
                    if new_hex:
                        fc.rgb = RGBColor.from_string(new_hex)
                        cnt += 1
                except Exception:
                    pass
        return cnt

    changed = 0
    if getattr(shape, "has_text_frame", False):
        changed += _recolor_frame(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                try:
                    changed += _recolor_frame(cell.text_frame)
                except Exception:
                    pass
    return changed


def update_theme_colors(pptx_path: Path, target_hex: str | None, palette: dict | None = None) -> None:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    if palette:
        variant_map = {
            "dk1":     palette.get("bodyText", target_hex or "1D2435"),
            "lt1":     palette.get("canvas", "FFFFFF"),
            "accent1": palette.get("accentA", target_hex or "000000"),
            "accent2": palette.get("accentB", target_hex or "000000"),
            "accent3": palette.get("sectionBg", "EEEEEE"),
            "accent4": palette.get("tableAltBg", "F5F5F5"),
            "accent5": palette.get("mutedText", "666666"),
            "accent6": palette.get("border", "DDDDDD"),
            "hlink":   palette.get("accentB", target_hex or "0563C1"),
        }
    elif target_hex:
        variant_map = {
            "accent1": target_hex,
            "accent2": shift_lightness(target_hex, 0.12),
            "accent3": shift_lightness(target_hex, -0.10),
            "accent4": shift_lightness(target_hex, 0.2),
            "accent5": shift_lightness(target_hex, -0.18),
            "accent6": shift_lightness(target_hex, 0.04),
            "hlink":   target_hex,
            "folHlink": shift_lightness(target_hex, -0.12),
        }
    else:
        return

    with zipfile.ZipFile(pptx_path, "r") as zin:
        entries: dict[str, bytes] = {}
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/") and item.filename.endswith(".xml"):
                root = etree.fromstring(data)
                updated = False
                for name, color_hex in variant_map.items():
                    node = root.find(f".//a:clrScheme/a:{name}", ns)
                    if node is None:
                        continue
                    for child in list(node):
                        node.remove(child)
                    srgb = etree.SubElement(
                        node,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
                    )
                    srgb.set("val", color_hex)
                    updated = True
                if updated:
                    data = etree.tostring(
                        root, xml_declaration=True, encoding="UTF-8", standalone="yes"
                    )
            entries[item.filename] = data

    with zipfile.ZipFile(pptx_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for filename, data in entries.items():
            zout.writestr(filename, data)


def find_shape_by_text(slide, near_text: str):
    """near_text を含む Shape を返す（見つからなければ None）。"""
    near_lower = near_text.lower()
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            if near_lower in shape.text_frame.text.lower():
                return shape
    return None


def insert_image(
    slide,
    image_path: str,
    position: str,
    width_pct: float,
    slide_width: int,
    slide_height: int,
    near_text: str = "",
    anchor_side: str = "right",
) -> bool:
    """DALL-E で生成した画像をスライドに挿入する。
    near_text が指定されていればその Shape の隣に配置し、見つからなければ固定 position を使う。
    """
    try:
        # widthPct をクランプ（5〜50%）
        width_pct = max(5.0, min(50.0, width_pct))
        width = int(slide_width * width_pct / 100)
        height = width  # アイコンは正方形
        margin = int(slide_width * 0.015)

        # anchorSide バリデーション
        if anchor_side not in {"left", "right", "above", "below"}:
            print(f"[edit_pptx] invalid anchorSide '{anchor_side}', defaulting to 'right'", file=sys.stderr)
            anchor_side = "right"

        # ヘッダー下限：スライド高さの20%以上に強制（タイトルバーへの被りを防ぐ）
        header_bottom = int(slide_height * 0.20)

        # nearText が指定されていれば Shape 相対配置を試みる
        if near_text:
            anchor_shape = find_shape_by_text(slide, near_text)
            if anchor_shape:
                s = anchor_shape

                # 左端シェイプの左に置こうとすると見切れる → 右に自動flip
                if anchor_side == "left" and s.left < width + margin * 2:
                    anchor_side = "right"

                if anchor_side == "right":
                    left = s.left + s.width + margin
                    top = s.top + (s.height - height) // 2
                elif anchor_side == "left":
                    left = s.left - width - margin
                    top = s.top + (s.height - height) // 2
                elif anchor_side == "above":
                    left = s.left + (s.width - width) // 2
                    top = s.top - height - margin
                else:  # below
                    left = s.left + (s.width - width) // 2
                    top = s.top + s.height + margin

                # ヘッダー回避 + スライド範囲内クランプ
                left = max(0, min(slide_width - width, left))
                top = max(header_bottom, min(slide_height - height, top))
                slide.shapes.add_picture(image_path, left, top, width, height)
                return True
            else:
                print(f"[edit_pptx] nearText '{near_text}' not found, falling back to position", file=sys.stderr)

        # 固定位置（フォールバック）：top系はヘッダー下から配置
        valid_positions = {"top-right", "top-left", "bottom-right", "bottom-left", "center"}
        if position not in valid_positions:
            position = "top-right"

        pos_map = {
            "top-right":    (slide_width - width - margin, header_bottom),
            "top-left":     (margin, header_bottom),
            "bottom-right": (slide_width - width - margin, slide_height - height - margin),
            "bottom-left":  (margin, slide_height - height - margin),
            "center":       ((slide_width - width) // 2, (slide_height - height) // 2),
        }
        left, top = pos_map[position]
        slide.shapes.add_picture(image_path, left, top, width, height)
        return True
    except Exception as e:
        print(f"[edit_pptx] insert_image failed: {e}", file=sys.stderr)
        return False


def add_bullets_to_shape(
    shape, add_bullets: list[dict], inserted_set: set | None = None
) -> bool:
    """afterText で特定した段落の直後に箇条書きを挿入する。書式は既存段落を XML コピーして保持。
    afterText が指定されていない・見つからない場合はその shape をスキップ（誤挿入防止）。
    inserted_set: 挿入済み afterText を追跡するセット。渡すと同一 afterText の二重挿入を防ぐ。"""
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    changed = False

    for item in add_bullets:
        after_text = (item.get("afterText") or "").strip()
        texts = [str(t)[:60] for t in (item.get("texts") or [])[:3] if str(t).strip()]
        if not texts:
            continue

        # afterText を含む段落を特定（para.runs は r_lst 未登録要素でクラッシュするため lxml 直アクセス）
        target_p = None
        for para in tf.paragraphs:
            para_text = "".join(
                t.text or "" for t in para._p.findall(f".//{{{NS}}}r/{{{NS}}}t")
            ).strip()
            if after_text and after_text in para_text:
                target_p = para._p
                break
        # afterText が空の場合のみ最終段落にフォールバック
        if target_p is None and not after_text and tf.paragraphs:
            target_p = tf.paragraphs[-1]._p
        if target_p is None:
            continue  # この shape には afterText が存在しない

        # 挿入済みとしてマーク（呼び出し元が次 shape でスキップできるよう）
        if inserted_set is not None:
            inserted_set.add(after_text)

        txBody = target_p.getparent()
        insert_idx = list(txBody).index(target_p) + 1

        for text in texts:
            # parse_xml で pptx 登録クラスとして生成（etree.fromstring は r_lst 未登録になる）
            new_p = parse_xml(etree.tostring(target_p))
            runs_in_new = new_p.findall(f"{{{NS}}}r")
            if runs_in_new:
                t_elem = runs_in_new[0].find(f"{{{NS}}}t")
                if t_elem is not None:
                    t_elem.text = text
                for extra in runs_in_new[1:]:
                    new_p.remove(extra)
            else:
                new_r = etree.SubElement(new_p, f"{{{NS}}}r")
                new_t = etree.SubElement(new_r, f"{{{NS}}}t")
                new_t.text = text
            txBody.insert(insert_idx, new_p)
            insert_idx += 1
            changed = True

    return changed


def _get_shape_y(shape, ns_a: str = "http://schemas.openxmlformats.org/drawingml/2006/main") -> int:
    """shape の Y 座標を EMU で返す（xfrm 優先、フォールバックは .top）"""
    xfrm = shape._element.find(f".//{{{ns_a}}}xfrm")
    if xfrm is not None:
        off = xfrm.find(f"{{{ns_a}}}off")
        if off is not None:
            try:
                return int(off.get("y", shape.top))
            except (ValueError, TypeError):
                pass
    return int(shape.top)


def _set_shape_y(sp_elem, new_y: int, ns_a: str = "http://schemas.openxmlformats.org/drawingml/2006/main") -> None:
    xfrm = sp_elem.find(f".//{{{ns_a}}}xfrm")
    if xfrm is None:
        return
    off = xfrm.find(f"{{{ns_a}}}off")
    if off is not None:
        off.set("y", str(int(new_y)))


def _cluster_shapes_by_y_band(shapes, gap_emu: int = 30_000) -> list[dict]:
    """
    Y座標バンドでshapeをクラスタリングする。
    gap_emu 以下の隙間は同一グループとみなす。
    各クラスタは dict(shapes=[...], y_range=(top, bot)) で返す。
    """
    if not shapes:
        return []
    sorted_s = sorted(shapes, key=_get_shape_y)
    clusters: list[dict] = []
    for shape in sorted_s:
        s_top = _get_shape_y(shape)
        s_bot = s_top + shape.height
        placed = False
        for cl in clusters:
            c_top, c_bot = cl["y_range"]
            if s_top <= c_bot + gap_emu and s_bot >= c_top - gap_emu:
                cl["shapes"].append(shape)
                cl["y_range"] = (min(c_top, s_top), max(c_bot, s_bot))
                placed = True
                break
        if not placed:
            clusters.append({"shapes": [shape], "y_range": (s_top, s_bot)})
    return clusters


def enable_text_to_fit_shape(shape) -> bool:
    """テキストフレームを TEXT_TO_FIT_SHAPE（normAutofit）にして自動フォント縮小を有効化する。
    XML の <a:bodyPr> 内の autofit 子要素を <a:normAutofit /> に差し替える。"""
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        NS_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(f"{{{NS_DML}}}bodyPr")
        if bodyPr is None:
            return False
        # 既存の autofit 要素（noAutofit / spAutoFit / normAutofit）を除去
        for tag_name in ("noAutofit", "spAutoFit", "normAutofit"):
            for el in bodyPr.findall(f"{{{NS_DML}}}{tag_name}"):
                bodyPr.remove(el)
        # <a:normAutofit /> を追加（TEXT_TO_FIT_SHAPE 相当）
        etree.SubElement(bodyPr, f"{{{NS_DML}}}normAutofit")
        return True
    except Exception as e:
        print(f"[enable_text_to_fit_shape] error on '{shape.name}': {e}", file=sys.stderr)
        return False


# ── テキスト増量後グループ再配置ヘルパー ──────────────────────────────────────────

def _get_shape_height(shape) -> int:
    """shape の高さを EMU で返す（xfrm 優先、フォールバックは .height）"""
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    xfrm = shape._element.find(f".//{{{NS_A}}}xfrm")
    if xfrm is not None:
        ext = xfrm.find(f"{{{NS_A}}}ext")
        if ext is not None:
            try:
                return int(ext.get("cy", shape.height))
            except (ValueError, TypeError):
                pass
    return int(shape.height)


def _set_shape_height(sp_elem, new_h: int) -> None:
    """shape._element の高さを EMU で設定する"""
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    xfrm = sp_elem.find(f".//{{{NS_A}}}xfrm")
    if xfrm is None:
        return
    ext = xfrm.find(f"{{{NS_A}}}ext")
    if ext is not None:
        ext.set("cy", str(int(new_h)))


def get_dominant_font_pt(shape) -> float:
    """shape 内で最も文字数が多い run のフォントサイズ（pt）を返す。未設定なら 18.0"""
    if not getattr(shape, "has_text_frame", False):
        return 18.0
    best_pt: float | None = None
    best_len = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size and len(run.text) > best_len:
                best_pt = run.font.size / 12700
                best_len = len(run.text)
    return best_pt if best_pt is not None else 18.0


def set_all_run_font_pt(shape, pt: float) -> bool:
    """shape 内の全 run のフォントサイズを pt に設定する。変更があれば True を返す"""
    if not getattr(shape, "has_text_frame", False):
        return False
    changed = False
    new_size_emu = int(pt * 12700)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            current = run.font.size
            if current is None or abs(current - new_size_emu) > 6350:
                run.font.size = new_size_emu
                changed = True
    return changed


def estimate_text_height_pt(text: str, box_w_pt: float, font_pt: float) -> float:
    """Meiryo 全角想定で text が box_w_pt 幅に収まるために必要な高さ（pt）を推定する"""
    chars_per_line = max(1.0, box_w_pt / (font_pt * 0.7))
    lines = math.ceil(len(text) / chars_per_line)
    return max(font_pt * 1.4, lines * font_pt * 1.4)


def reflow_group_after_text_change(
    slide,
    changed_shape_name: str,
    slide_height_emu: int,
    min_font_pt: float = 9.0,
) -> str | None:
    """
    replaceText / appendToRun 後に shape 高さ / フォント / 隣接グループ押し下げを行う。
    プリフライトで収まらないと判定した場合は shape を変更せず警告文字列を返す。
    正常終了なら None を返す。
    """
    # 変更 shape を探す（グループ shape の子も含む）
    shapes_by_name: dict[str, Any] = {}
    for s in slide.shapes:
        shapes_by_name[s.name] = s
        for child in getattr(s, "shapes", []):
            shapes_by_name[child.name] = child
    changed_shape = shapes_by_name.get(changed_shape_name)
    if changed_shape is None or not getattr(changed_shape, "has_text_frame", False):
        return None

    # Step 1: normAutofit を設定（常に安全）
    enable_text_to_fit_shape(changed_shape)

    text = (changed_shape.text_frame.text or "").strip()
    if not text:
        return None

    shape_w_pt = max(1.0, changed_shape.width / 12700)
    shape_h_pt = max(1.0, _get_shape_height(changed_shape) / 12700)
    orig_font_pt = get_dominant_font_pt(changed_shape)
    needed_h_pt = estimate_text_height_pt(text, shape_w_pt, orig_font_pt)

    # normAutofit で吸収可能なら終了
    if needed_h_pt <= shape_h_pt * 1.1:
        return None

    # Step 2: フォント縮小（floor: min_font_pt or 0.70 × 元サイズ・常に安全）
    floor_pt = max(min_font_pt, orig_font_pt * 0.70)
    target_pt = orig_font_pt - 1.0
    while target_pt >= floor_pt:
        needed_h_pt = estimate_text_height_pt(text, shape_w_pt, target_pt)
        if needed_h_pt <= shape_h_pt * 1.05:
            set_all_run_font_pt(changed_shape, target_pt)
            print(
                f"[reflow_group] font shrunk: '{changed_shape_name}' "
                f"{orig_font_pt:.0f}pt -> {target_pt:.0f}pt",
                file=sys.stderr,
            )
            return None
        target_pt -= 1.0

    # floor_pt でフォントを固定（変更は安全）
    target_pt = floor_pt
    if target_pt < orig_font_pt:
        set_all_run_font_pt(changed_shape, target_pt)
    needed_h_pt = estimate_text_height_pt(text, shape_w_pt, target_pt)

    # 拡張量を計算
    extra_pt = max(0.0, needed_h_pt - shape_h_pt)
    if extra_pt < 1.0:
        return None
    extra_h_emu = int(extra_pt * 12700)

    old_shape_h = _get_shape_height(changed_shape)
    old_shape_bottom_emu = _get_shape_y(changed_shape) + old_shape_h

    # Step 3: クラスタリング（拡張前に行い old_cluster_bottom_emu を確定）
    content_shapes = [
        s for s in slide.shapes
        if _get_shape_height(s) <= slide_height_emu * 0.6
    ]
    all_clusters = _cluster_shapes_by_y_band(content_shapes, gap_emu=30_000)
    changed_cluster = None
    old_cluster_bottom_emu = old_shape_bottom_emu  # フォールバック
    for cl in all_clusters:
        if any(s.name == changed_shape_name for s in cl["shapes"]):
            changed_cluster = cl
            old_cluster_bottom_emu = cl["y_range"][1]  # 拡張前の底辺
            break

    # Step 4: プリフライト（shape を一切変更せず溢れるか確認）
    post_exp_max = 0
    for cl in all_clusters:
        if cl is changed_cluster:
            post_exp_max = max(post_exp_max, cl["y_range"][1] + extra_h_emu)
        elif cl["y_range"][0] >= old_cluster_bottom_emu - 20_000:
            # 押し下げ対象クラスタ
            post_exp_max = max(post_exp_max, cl["y_range"][1] + extra_h_emu)
        else:
            post_exp_max = max(post_exp_max, cl["y_range"][1])

    if post_exp_max > slide_height_emu:
        overflow_pt = (post_exp_max - slide_height_emu) / 12700
        return (
            f"テキスト追加後の高さ調整でスライド下端から約{overflow_pt:.0f}ptはみ出します。"
            f"文字量を減らすか、スライドを分割してください。"
        )

    # プリフライト通過 → 実際に拡張する
    _set_shape_height(changed_shape._element, old_shape_h + extra_h_emu)
    print(
        f"[reflow_group] shape height expanded: '{changed_shape_name}' "
        f"+{extra_h_emu}EMU ({extra_pt:.1f}pt)",
        file=sys.stderr,
    )

    if changed_cluster is not None:
        changed_left = int(changed_shape.left)
        changed_right = changed_left + int(changed_shape.width)

        # Step 4a: 背景カード shape を拡張（画像・アイコンは除外）
        for grp_shape in changed_cluster["shapes"]:
            if grp_shape.name == changed_shape_name:
                continue
            no_text = (
                not getattr(grp_shape, "has_text_frame", False)
                or not (grp_shape.text_frame.text or "").strip()
            )
            not_picture = grp_shape.shape_type != MSO_SHAPE_TYPE.PICTURE
            # テキスト shape を水平方向にカバーする shape のみ対象
            h_covers = (
                int(grp_shape.left) <= changed_left + 50_000
                and int(grp_shape.left) + int(grp_shape.width) >= changed_right - 50_000
            )
            grp_bottom = _get_shape_y(grp_shape) + _get_shape_height(grp_shape)
            if no_text and not_picture and h_covers and grp_bottom >= old_shape_bottom_emu - 50_000:
                _set_shape_height(grp_shape._element, _get_shape_height(grp_shape) + extra_h_emu)
                print(
                    f"[reflow_group] bg shape expanded: '{grp_shape.name}' +{extra_h_emu}EMU",
                    file=sys.stderr,
                )

        # Step 4b: 下方クラスタを押し下げ（old_cluster_bottom_emu を基準にする）
        for cl in all_clusters:
            if cl is changed_cluster:
                continue
            if cl["y_range"][0] >= old_cluster_bottom_emu - 20_000:
                for s in cl["shapes"]:
                    _set_shape_y(s._element, _get_shape_y(s) + extra_h_emu)
                print(
                    f"[reflow_group] cluster pushed: top={cl['y_range'][0]} +{extra_h_emu}EMU",
                    file=sys.stderr,
                )

    return None


def detect_text_overflow_candidates(slide, slide_height: int) -> list[str]:
    """
    テキストが shape に収まらない可能性がある shape 名のリストを返す。
    Meiryo 実描画を想定した粗い推定（FP あり）。
    """
    overflow: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if len(text) < 10:
            continue
        min_pt = 18
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size // 12700
                        if 6 <= pt < min_pt:
                            min_pt = pt
        except Exception:
            pass
        w_pt = max(1.0, shape.width / 12700)
        h_pt = max(1.0, shape.height / 12700)
        # Meiryo 全角: 1文字 ≈ font_pt * 0.7 pt 幅、行高 ≈ font_pt * 1.4 pt
        chars_per_line = w_pt / (min_pt * 0.7)
        lines_avail = h_pt / (min_pt * 1.4)
        capacity = chars_per_line * lines_avail
        if len(text) > capacity * 1.3:
            overflow.append(shape.name)
    return overflow


def copy_shape_block(slide, action: dict, slide_height: int) -> bool:
    """
    アイテムグループ（背景カード＋アイコン＋テキスト）単位で shape をコピーし追加する。
    Y座標バンドクラスタリングで、テキスト枠だけでなくカード・アイコン画像も含めて複製する。
    """
    from pptx.oxml.ns import qn

    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    heading_name = (action.get("headingShapeName") or "").strip()
    desc_name    = (action.get("descShapeName")    or "").strip()
    heading_text = (action.get("headingText")      or "").strip()
    desc_text    = (action.get("descText")         or "").strip()
    group_shape_names: list[str] = action.get("groupShapeNames") or []

    if not (heading_name and desc_name):
        print("[copy_shape_block] missing shape names", file=sys.stderr)
        return False

    shapes_by_name = {s.name: s for s in slide.shapes}
    src_heading = shapes_by_name.get(heading_name)
    src_desc    = shapes_by_name.get(desc_name)
    if not src_heading or not src_desc:
        print(f"[copy_shape_block] shapes not found: {heading_name!r}, {desc_name!r}", file=sys.stderr)
        return False

    def set_text(sp_elem, text: str) -> None:
        for t in sp_elem.findall(f".//{{{NS_A}}}t"):
            t.text = ""
        t_elems = sp_elem.findall(f".//{{{NS_A}}}t")
        if t_elems:
            t_elems[0].text = text

    # ── アンカー shape の特定（LLM 指定のテキスト shape 名、または heading+desc のみ）──
    anchor_names_set = set(group_shape_names) if group_shape_names else {heading_name, desc_name}
    # heading/desc は LLM が groupShapeNames に含め忘れた場合でも常にアンカーとして保証する
    anchor_names_set |= {heading_name, desc_name}
    anchor_shapes = [shapes_by_name[n] for n in anchor_names_set if n in shapes_by_name]
    if not anchor_shapes:
        anchor_shapes = [src_heading, src_desc]

    # ── 背景大型 shape を除外してコンテンツ shape だけを対象にする ──
    # slide_height の 60% 以上の height は装飾・背景として除外
    content_shapes = [s for s in slide.shapes if s.height <= slide_height * 0.6]

    # ── Y帯クラスタリング（全コンテンツ shape）──
    all_clusters = _cluster_shapes_by_y_band(content_shapes, gap_emu=30_000)

    # ── アンカーを含むクラスタ = アイテムグループ ──
    item_groups = [
        cl for cl in all_clusters
        if any(s.name in anchor_names_set for s in cl["shapes"])
    ]
    item_groups.sort(key=lambda cl: cl["y_range"][0])

    if not item_groups:
        # フォールバック: src_heading + src_desc の2 shape のみ
        top0 = _get_shape_y(src_heading)
        bot0 = _get_shape_y(src_desc) + src_desc.height
        item_groups = [{"shapes": [src_heading, src_desc], "y_range": (top0, bot0)}]

    # ── スペース計算 ──
    first_y          = item_groups[0]["y_range"][0]
    available        = slide_height - first_y
    group_heights    = [cl["y_range"][1] - cl["y_range"][0] for cl in item_groups]
    existing_total_h = sum(group_heights)
    template_group   = item_groups[-1]
    template_h       = template_group["y_range"][1] - template_group["y_range"][0]
    n_total          = len(item_groups) + 1
    total_h          = existing_total_h + template_h

    if total_h > available:
        print(
            f"[copy_shape_block] not enough space: need={total_h} available={available} "
            f"groups={len(item_groups)} shapes_per_group={len(template_group['shapes'])}",
            file=sys.stderr,
        )
        return False

    gap = int((available - total_h) / max(n_total - 1, 1))

    # ── 既存グループの Y 再配置（テキスト・カード・アイコンを全て移動）──
    cur_y = first_y
    actual_template_top: int | None = None
    for grp in item_groups:
        # template_group が移動される直前の cur_y を記録（移動後の相対 offset 計算に使用）
        if grp is template_group:
            actual_template_top = cur_y
        grp_top = grp["y_range"][0]
        delta   = cur_y - grp_top
        if delta != 0:
            for shape in grp["shapes"]:
                _set_shape_y(shape._element, _get_shape_y(shape) + delta, NS_A)
        cur_y += (grp["y_range"][1] - grp_top) + gap

    new_group_top = cur_y
    # actual_template_top: 移動後のテンプレートグループ先頭 Y（y_range[0] は移動前の古い値のため使用不可）
    template_top  = actual_template_top if actual_template_top is not None else template_group["y_range"][0]

    # ── 新グループのコピー（テンプレートグループの全 shape を複製）──
    existing_ids: set[int] = set()
    for elem in slide.shapes._spTree.iter():
        id_val = elem.get("id")
        if id_val is not None:
            try:
                existing_ids.add(int(id_val))
            except ValueError:
                pass
    max_id = max(existing_ids, default=100)

    copied = 0
    for i, src in enumerate(sorted(template_group["shapes"], key=_get_shape_y)):
        offset_y = _get_shape_y(src) - template_top
        new_y    = new_group_top + offset_y

        new_elem = deepcopy(src._element)
        cNvPr = new_elem.find(f".//{qn('p:cNvPr')}")
        if cNvPr is not None:
            cNvPr.set("id",   str(max_id + i + 1))
            cNvPr.set("name", cNvPr.get("name", "") + "_new")
        _set_shape_y(new_elem, new_y, NS_A)

        # 見出し・説明テキストのみ更新（背景カード・アイコンはそのまま複製）
        if src.name == heading_name and heading_text:
            set_text(new_elem, heading_text)
        elif src.name == desc_name and desc_text:
            set_text(new_elem, desc_text)

        slide.shapes._spTree.append(new_elem)
        copied += 1

    print(
        f"[copy_shape_block] added group at y={new_group_top} gap={gap} "
        f"(item_groups={len(item_groups)}, copied_shapes={copied}, "
        f"template_h={template_h}, available={available})",
        file=sys.stderr,
    )
    return True


def repeat_copy_shape_block(slide, action: dict, slide_height: int) -> dict:
    """
    「項目数をNにする」用。
    targetItemCount 指定時: 一括レイアウト計算（upward shift + gap圧縮）で全グループを配置。
    targetItemCount 未指定時: copy_shape_block を逐次呼び出す従来モード。
    Returns dict: {added, required_add, final_count, target_count, success}
    """
    from pptx.oxml.ns import qn
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    heading_name  = (action.get("headingShapeName") or "").strip()
    desc_name     = (action.get("descShapeName")    or "").strip()
    new_items     = list(action.get("newItems") or [])
    initial_group = list(action.get("groupShapeNames") or [])
    target_count  = action.get("targetItemCount")

    def _result(added=0, required_add=0, final_count=None, success=False, layout_mode=None):
        r: dict = {"added": added, "required_add": required_add,
                   "final_count": final_count, "target_count": target_count, "success": success}
        if layout_mode:
            r["layout_mode"] = layout_mode
        return r

    if not (heading_name and desc_name) or not new_items:
        print("[repeat_copy_shape_block] missing headingShapeName/descShapeName or empty newItems",
              file=sys.stderr)
        return _result(required_add=len(new_items))

    # アンカーセット + 現在グループ検出（LLM currentCount 自己補正にも使用）
    anchor_set = set(initial_group) | {heading_name, desc_name}
    all_content = [s for s in slide.shapes if s.height <= slide_height * 0.6]
    all_cl = _cluster_shapes_by_y_band(all_content, gap_emu=30_000)

    # anchor shapes を含むクラスタを起点として検出
    anchor_clusters = sorted(
        [cl for cl in all_cl if any(s.name in anchor_set for s in cl["shapes"])],
        key=lambda cl: cl["y_range"][0],
    )
    # LLM が groupShapeNames を部分的にしか返さない場合の補完:
    # 最初〜最後 anchor クラスタの間にあり、かつ同程度の高さのクラスタも item 群に含める
    if anchor_clusters:
        tpl_h_est = anchor_clusters[-1]["y_range"][1] - anchor_clusters[-1]["y_range"][0]
        h_lo, h_hi = tpl_h_est * 0.6, tpl_h_est * 1.4
        anchor_ids   = {id(cl) for cl in anchor_clusters}
        first_anchor_y = anchor_clusters[0]["y_range"][0]
        last_anchor_y  = anchor_clusters[-1]["y_range"][1]
        between = [
            cl for cl in all_cl
            if id(cl) not in anchor_ids
            and cl["y_range"][0] >= first_anchor_y - 30_000
            and cl["y_range"][1] <= last_anchor_y + 30_000
            and h_lo <= (cl["y_range"][1] - cl["y_range"][0]) <= h_hi
        ]
        item_clusters = sorted(anchor_clusters + between, key=lambda cl: cl["y_range"][0])
    else:
        item_clusters = []
    actual_current = len(item_clusters)

    if isinstance(target_count, int) and target_count > 0:
        required_add = max(0, target_count - actual_current)
        if required_add == 0:
            print(f"[repeat_copy_shape_block] already at target ({actual_current}), skip", file=sys.stderr)
            return _result(added=0, required_add=0, final_count=actual_current, success=True)
        print(
            f"[repeat_copy_shape_block] targetCount={target_count} detected={actual_current} "
            f"requiredAdd={required_add}",
            file=sys.stderr,
        )
        if required_add < len(new_items):
            new_items = new_items[:required_add]
        elif required_add > len(new_items):
            # newItems が足りない → 部分編集を防ぐため変更なしで即失敗返却
            print(
                f"[repeat_copy_shape_block] ERROR: need {required_add} but only "
                f"{len(new_items)} items provided → abort (no partial edit)",
                file=sys.stderr,
            )
            return _result(added=0, required_add=required_add,
                           final_count=actual_current, success=False)
    else:
        required_add = len(new_items)

    if not item_clusters:
        print("[repeat_copy_shape_block] no item clusters found, abort", file=sys.stderr)
        return _result(required_add=required_add)

    template_cluster = item_clusters[-1]
    template_h = template_cluster["y_range"][1] - template_cluster["y_range"][0]
    n_existing = actual_current
    n_new = len(new_items)
    n_total = n_existing + n_new
    first_y = item_clusters[0]["y_range"][0]

    # ── BATCH LAYOUT MODE（targetItemCount 指定時）────────────────────────────
    if isinstance(target_count, int) and target_count > 0:
        # レイアウト定数
        BOTTOM_MARGIN = 300_000   # ~8mm: スライド下端との最小距離
        TOP_MARGIN    = 150_000   # ~4mm: タイトル下端からの最小距離
        MAX_GAP       = 350_000   # ~1cm: グループ間隔の上限（広がりすぎ防止）
        MIN_GAP       =  50_000   # ~1.4mm: グループ間隔の下限

        # タイトル等の非アイテムクラスタから usable_top を決定
        non_item_above = [
            cl for cl in all_cl
            if not any(s.name in anchor_set for s in cl["shapes"])
            and cl["y_range"][1] <= first_y
        ]
        if non_item_above:
            title_bot = max(cl["y_range"][1] for cl in non_item_above)
            usable_top = title_bot + TOP_MARGIN
        else:
            usable_top = first_y

        usable_bottom   = slide_height - BOTTOM_MARGIN
        usable_h        = max(0, usable_bottom - usable_top)
        # 各クラスタの実際の高さを使う（template_h で一律計算すると背高クラスタが誤差を生む）
        existing_heights = [cl["y_range"][1] - cl["y_range"][0] for cl in item_clusters]
        total_content_h  = sum(existing_heights) + template_h * n_new

        def _try_layout(avail_h: int):
            """avail_h に収まる (gap, block_h) を返す。収まらなければ None。"""
            min_block = total_content_h + MIN_GAP * (n_total - 1)
            if min_block > avail_h:
                return None
            raw_gap = (avail_h - total_content_h) / (n_total - 1) if n_total > 1 else 0
            g  = int(min(MAX_GAP, max(MIN_GAP, raw_gap))) if n_total > 1 else 0
            bh = total_content_h + g * (n_total - 1)
            return g, bh

        # 段階1: BOTTOM_MARGIN 確保・MAX_GAP 上限・縦中央寄せ
        lo = _try_layout(usable_h)
        if lo:
            gap, block_h   = lo
            layout_first_y = usable_top + max(0, (usable_h - block_h) // 2)
            layout_mode    = "balanced"
        else:
            # 段階2: BOTTOM_MARGIN を犠牲にした tight 配置
            tight_h = slide_height - usable_top
            lo = _try_layout(tight_h)
            if lo:
                gap, block_h   = lo
                layout_first_y = usable_top + max(0, (tight_h - block_h) // 2)
                layout_mode    = "tight"
            else:
                print(
                    f"[repeat_copy_shape_block] FAIL: need={total_content_h} "
                    f"usable_h={usable_h} tight_h={tight_h} n={n_total} h={template_h}",
                    file=sys.stderr,
                )
                print(
                    f"[repeat_copy_shape_block] total added=0/{n_new} "
                    f"finalCount={actual_current} success=false",
                    file=sys.stderr,
                )
                return _result(added=0, required_add=required_add,
                               final_count=actual_current, success=False)

        bottom_margin_actual = slide_height - (layout_first_y + block_h)
        print(
            f"[repeat_copy_shape_block] layout mode={layout_mode} "
            f"usableTop={usable_top} usableBottom={usable_bottom} "
            f"blockH={block_h} firstY={layout_first_y} gap={gap} "
            f"bottomMargin={bottom_margin_actual}",
            file=sys.stderr,
        )

        # 既存グループを layout_first_y から順に再配置（各自の実高さでステップ）
        cur_y = layout_first_y
        new_template_top = None
        for i, grp in enumerate(item_clusters):
            if i == len(item_clusters) - 1:
                new_template_top = cur_y
            delta = cur_y - grp["y_range"][0]
            if delta != 0:
                for shape in grp["shapes"]:
                    _set_shape_y(shape._element, _get_shape_y(shape) + delta, NS_A)
            cur_y += existing_heights[i] + gap  # 実高さでステップ（template_h ではない）

        if new_template_top is None:
            new_template_top = layout_first_y

        # max shape ID 取得
        existing_ids: set[int] = set()
        for elem in slide.shapes._spTree.iter():
            id_val = elem.get("id")
            if id_val:
                try:
                    existing_ids.add(int(id_val))
                except ValueError:
                    pass
        max_id = max(existing_ids, default=100)

        def _set_text(sp_elem, text: str) -> None:
            for t in sp_elem.findall(f".//{{{NS_A}}}t"):
                t.text = ""
            t_elems = sp_elem.findall(f".//{{{NS_A}}}t")
            if t_elems:
                t_elems[0].text = text

        template_shapes = sorted(template_cluster["shapes"], key=_get_shape_y)
        id_offset = 0

        for item_idx, item in enumerate(new_items):
            new_group_top = cur_y
            heading_text = (item.get("headingText") or "").strip()
            desc_text    = (item.get("descText")    or "").strip()

            for i, src in enumerate(template_shapes):
                offset_y = _get_shape_y(src) - new_template_top
                new_y    = new_group_top + offset_y

                new_elem = deepcopy(src._element)
                cNvPr = new_elem.find(f".//{qn('p:cNvPr')}")
                if cNvPr is not None:
                    cNvPr.set("id",   str(max_id + id_offset + i + 1))
                    cNvPr.set("name", cNvPr.get("name", "") + f"_new{item_idx + 1}")
                _set_shape_y(new_elem, new_y, NS_A)

                if src.name == heading_name and src.name == desc_name:
                    # 同一shape: 見出しと説明を改行で結合して設定
                    combined = "\n".join(t for t in [heading_text, desc_text] if t)
                    if combined:
                        _set_text(new_elem, combined)
                elif src.name == heading_name and heading_text:
                    _set_text(new_elem, heading_text)
                elif src.name == desc_name and desc_text:
                    _set_text(new_elem, desc_text)

                slide.shapes._spTree.append(new_elem)

            id_offset += len(template_shapes)
            cur_y += template_h + gap
            print(f"[repeat_copy_shape_block] item {item_idx + 1} added", file=sys.stderr)

        added = len(new_items)
        final_count = n_total
        print(
            f"[repeat_copy_shape_block] total added={added}/{required_add} "
            f"finalCount={final_count} bottomMargin={bottom_margin_actual} success=true",
            file=sys.stderr,
        )
        return _result(added=added, required_add=required_add,
                       final_count=final_count, success=True, layout_mode=layout_mode)

    # ── ITERATIVE MODE（targetItemCount 未指定時）─────────────────────────────
    cur_heading = heading_name
    cur_desc    = desc_name
    accumulated_group = list(anchor_set)

    added = 0
    for idx, item in enumerate(new_items):
        names_before = {s.name for s in slide.shapes}

        single_action = {
            "headingShapeName": cur_heading,
            "descShapeName":    cur_desc,
            "groupShapeNames":  accumulated_group,
            "headingText":      (item.get("headingText") or "").strip(),
            "descText":         (item.get("descText")    or "").strip(),
        }
        ok = copy_shape_block(slide, single_action, slide_height)
        if not ok:
            print(f"[repeat_copy_shape_block] item {idx + 1}/{len(new_items)} failed → stop",
                  file=sys.stderr)
            break

        added += 1

        all_after = {s.name: s for s in slide.shapes}
        new_names = [n for n in all_after if n not in names_before]
        new_h = next((n for n in new_names if n.startswith(cur_heading)), None)
        new_d = next((n for n in new_names if n.startswith(cur_desc)),    None)

        if new_h and new_d:
            cur_heading = new_h
            cur_desc    = new_d
        else:
            cur_heading = cur_heading + "_new"
            cur_desc    = cur_desc    + "_new"

        accumulated_group = accumulated_group + [cur_heading, cur_desc]
        print(f"[repeat_copy_shape_block] item {idx + 1} added → next={cur_heading}/{cur_desc}",
              file=sys.stderr)

    print(f"[repeat_copy_shape_block] total added={added}/{len(new_items)}", file=sys.stderr)
    success = added >= required_add
    return _result(added=added, required_add=required_add,
                   final_count=actual_current + added, success=success)


def count_all_run_chars(prs) -> int:
    """全スライドの全run文字数合計を返す（charsBefore/charsAfter計算用）"""
    total = 0
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        total += len(run.text)
    return total


def _is_title_like_shape(shape, slide_height: int) -> bool:
    try:
        if getattr(shape, "is_placeholder", False):
            ph_type = str(shape.placeholder_format.type).lower()
            if "title" in ph_type:
                return True
    except Exception:
        pass
    try:
        text = (shape.text or "").strip()
        return bool(text) and int(shape.top) < slide_height * 0.18 and int(shape.height) < slide_height * 0.18
    except Exception:
        return False


def _remove_body_content_shapes(slide, slide_height: int) -> None:
    for shape in list(slide.shapes):
        try:
            if _is_title_like_shape(shape, slide_height):
                continue
            if int(shape.top) < slide_height * 0.16 and int(shape.height) < slide_height * 0.14:
                continue
            is_text = bool(getattr(shape, "has_text_frame", False) and (shape.text or "").strip())
            is_auto = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            if is_text or is_auto or is_picture:
                slide.shapes._spTree.remove(shape.element)
        except Exception:
            continue


def _set_textbox_text(shape, heading: str, body: str, card_h: int = 2200000, heading_color: str = "1F2937", body_color: str = "4B5563") -> None:
    scale = min(1.3, max(1.0, card_h / 2200000))
    heading_pt = round(17 * scale)
    body_pt = round(13 * scale)
    space_pt = round(8 * scale)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.margin_top = Emu(90000)
    tf.margin_bottom = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = heading
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(heading_pt)
        p.runs[0].font.color.rgb = RGBColor.from_string(heading_color)
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(space_pt)
        if p2.runs:
            p2.runs[0].font.size = Pt(body_pt)
            p2.runs[0].font.color.rgb = RGBColor.from_string(body_color)


def convert_slide_to_cards(slide, action: dict, slide_width: int, slide_height: int, deck_accent: str | None = None, palette: dict | None = None) -> bool:
    cards = action.get("cards") or []
    if not isinstance(cards, list):
        return False
    clean_cards: list[dict] = []
    for card in cards[:6]:
        if not isinstance(card, dict):
            continue
        heading = str(card.get("heading") or "").strip()
        body = str(card.get("body") or "").strip()
        if heading or body:
            clean_cards.append({"heading": heading[:40], "body": body[:180]})
    if not clean_cards:
        return False

    _remove_body_content_shapes(slide, slide_height)

    margin_x = int(slide_width * 0.07)
    content_top = int(slide_height * 0.23)
    content_bottom = int(slide_height * 0.90)
    gap_x = int(slide_width * 0.025)
    gap_y = int(slide_height * 0.035)
    count = len(clean_cards)
    cols = 1 if count == 1 else 2
    rows = math.ceil(count / cols)
    card_w = int((slide_width - margin_x * 2 - gap_x * (cols - 1)) / cols)
    card_h = int((content_bottom - content_top - gap_y * (rows - 1)) / rows)
    accent_colors = ["2563EB", "059669", "D97706", "7C3AED", "DC2626", "0891B2"]
    surface = palette.get("surface", "FFFFFF") if palette else "FFFFFF"
    border_color = palette.get("border", "D1D5DB") if palette else "D1D5DB"
    heading_color = palette.get("bodyText", "1F2937") if palette else "1F2937"
    body_color = palette.get("mutedText", "4B5563") if palette else "4B5563"
    palette_accent = palette.get("accentA") if palette else None

    for i, card in enumerate(clean_cards):
        row = i // cols
        col = i % cols
        left = margin_x + col * (card_w + gap_x)
        top = content_top + row * (card_h + gap_y)
        accent = deck_accent or palette_accent or accent_colors[i % len(accent_colors)]

        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Emu(left),
            Emu(top),
            Emu(card_w),
            Emu(card_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(surface)
        rect.line.color.rgb = RGBColor.from_string(border_color)
        rect.line.width = Pt(1.1)
        _set_textbox_text(rect, card["heading"], card["body"], card_h, heading_color, body_color)

        dot_size = int(min(card_w, card_h) * 0.12)
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Emu(left + int(card_w * 0.82)),
            Emu(top + int(card_h * 0.10)),
            Emu(dot_size),
            Emu(dot_size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(accent)
        dot.line.color.rgb = RGBColor.from_string(accent)

    return True


def _remap_relationships(sp_elem, ref_part, target_part) -> None:
    """Remap all r:embed / r:link rIds in sp_elem from ref_part to target_part.

    python-pptx stores image / hyperlink relationships per-slide in .rels files.
    After deepcopying a shape element, the embedded rIds still point to the
    source slide's relationships. This function re-registers each referenced
    part (or external URL) in the target slide's part and updates the XML attrs
    so the copied shape resolves correctly.
    """
    REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    embed_attr = f"{{{REL_NS}}}embed"
    link_attr  = f"{{{REL_NS}}}link"

    rId_cache: dict[str, str] = {}

    for elem in sp_elem.iter():
        for attr_name in (embed_attr, link_attr):
            old_rId = elem.get(attr_name)
            if not old_rId:
                continue
            if old_rId in rId_cache:
                elem.set(attr_name, rId_cache[old_rId])
                continue
            try:
                if old_rId not in ref_part.rels:
                    continue
                rel = ref_part.rels[old_rId]
                if rel.is_external:
                    new_rId = target_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    source_part = ref_part.related_part(old_rId)
                    new_rId = target_part.relate_to(source_part, rel.reltype)
                rId_cache[old_rId] = new_rId
                elem.set(attr_name, new_rId)
            except Exception as exc:
                print(
                    f"[copy_layout] WARNING: failed to remap rId={old_rId}: {exc}",
                    file=sys.stderr,
                )


def copy_slide_layout_from_reference(
    prs,
    target_slide_index: int,
    reference_slide_index: int,
    preserve_target_text: bool = True,
) -> bool:
    """Copy body shape layout/style from reference slide to target slide.

    Steps:
    1. Collect target slide's non-title text lines (if preserve_target_text).
    2. Remove all body/content shapes from target slide.
    3. Deep-copy content shapes from reference slide to target.
    4. Replace text in copied shapes with preserved target texts.
    """
    total = len(prs.slides)
    if not (0 <= target_slide_index < total) or not (0 <= reference_slide_index < total):
        return False
    if target_slide_index == reference_slide_index:
        return False

    target_slide = prs.slides[target_slide_index]
    ref_slide = prs.slides[reference_slide_index]
    slide_height = int(prs.slide_height)

    # 1. Collect target non-title texts in order
    target_texts: list[str] = []
    if preserve_target_text:
        for shape in list(iter_shapes(target_slide.shapes)):
            if _is_title_like_shape(shape, slide_height):
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    target_texts.append(t)

    # 2. Remove body/content shapes from target (title is kept by _remove_body_content_shapes)
    _remove_body_content_shapes(target_slide, slide_height)

    # 3. Deep-copy non-title shapes from reference to target.
    # r:embed / r:link rIds are per-slide, so re-register each referenced part
    # in target_slide.part before appending the element.
    copied_any = False
    for shape in list(iter_shapes(ref_slide.shapes)):
        if _is_title_like_shape(shape, slide_height):
            continue
        sp_elem = deepcopy(shape._element)
        _remap_relationships(sp_elem, ref_slide.part, target_slide.part)
        target_slide.shapes._spTree.append(sp_elem)
        copied_any = True

    if not copied_any:
        return False

    # 4. Replace text in copied body shapes with preserved target texts.
    #    参照元が「見出し + 説明」の繰り返しレイアウトで、ターゲット側の文章数が
    #    その項目数と一致する場合は、各文章を1項目ずつに整理して流し込む。
    #    従来の単純な順次流し込みでは、4文章が最初の2項目にだけ入り、残り2項目が
    #    空欄になるため、見出しは本文先頭から短く補完する。
    #    target_texts が尽きた後の残余シェイプは空にして参照元テキストを残さない。
    if preserve_target_text:
        body_shapes = [
            s for s in iter_shapes(target_slide.shapes)
            if not _is_title_like_shape(s, slide_height) and getattr(s, "has_text_frame", False)
        ]
        text_slots = [
            para
            for shape in body_shapes
            for para in shape.text_frame.paragraphs
            if para.runs and para.text.strip()
        ]

        def _heading_and_body(text: str) -> tuple[str, str]:
            cleaned = text.strip()
            for separator in ("：", ":", "。", "、"):
                pos = cleaned.find(separator)
                if 2 <= pos <= 24:
                    heading = cleaned[:pos].strip()
                    body = cleaned[pos + 1:].strip() or cleaned
                    return heading, body
            heading = cleaned[:20].strip()
            if len(cleaned) > 20:
                heading = heading.rstrip("、。・ ") + "…"
            return heading or "要点", cleaned

        # 見出し/説明の2枠 × N項目、かつターゲット文章がN件なら項目単位で配置する。
        paired_layout = (
            len(target_texts) >= 2
            and len(text_slots) >= len(target_texts) * 2
            and len(text_slots) % 2 == 0
        )
        replacement_texts: list[str]
        if paired_layout:
            replacement_texts = []
            for target_text in target_texts:
                heading, body = _heading_and_body(target_text)
                replacement_texts.extend([heading, body])
            print(
                f"[copy_layout] organized {len(target_texts)} texts into heading/body pairs",
                file=sys.stderr,
            )
        else:
            replacement_texts = target_texts

        for text_idx, para in enumerate(text_slots):
            if text_idx < len(replacement_texts):
                # Keep run formatting; replace text content only
                para.runs[0].text = replacement_texts[text_idx]
                for run in para.runs[1:]:
                    run.text = ""
            else:
                # target_texts 使い切り → 参照元テキストを残さないよう空にする
                for run in para.runs:
                    run.text = ""

    return True


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--plan", required=True)
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)
    plan_path = Path(args.plan)

    plan: dict[str, Any] = json.loads(plan_path.read_text(encoding="utf-8"))
    prs = Presentation(str(input_path))

    deck_edits = plan.get("deckEdits") or {}
    target_hex = normalize_hex(deck_edits.get("accentColor"))
    palette: dict | None = deck_edits.get("palette") if isinstance(deck_edits.get("palette"), dict) else None
    palette_key: str | None = deck_edits.get("paletteKey") or None
    effective_hex = target_hex or (palette.get("accentA") if palette else None)
    font_face = (deck_edits.get("fontFace") or "").strip() or None
    preserve_text_colors: bool = bool(deck_edits.get("preserveTextColors", True))

    slide_edit_map = {
        int(item.get("slideIndex")): item
        for item in (plan.get("slideEdits") or [])
        if item.get("slideIndex") is not None
    }

    # imageInserts: slideIndex → list of inserts (-1 = 全スライドに適用)
    image_insert_map: dict[int, list[dict]] = {}
    all_slides_inserts: list[dict] = []
    for item in (plan.get("imageInserts") or []):
        si = item.get("slideIndex")
        if si is not None:
            if int(si) == -1:
                all_slides_inserts.append(item)
            else:
                image_insert_map.setdefault(int(si), []).append(item)

    chars_before = count_all_run_chars(prs)
    changed_slides: set[int] = set()
    inserted_images: int = 0
    layout_warnings: list[str] = []
    changed_fills: int = 0
    changed_texts: int = 0
    changed_lines: int = 0
    item_count_slide_results: dict[int, dict] = {}

    # 範囲外 slideIndex の検出（②）
    total_slides = len(prs.slides)
    out_of_range_indices = sorted(
        si for si in slide_edit_map if si < 0 or si >= total_slides
    )
    if out_of_range_indices:
        print(
            f"[edit_pptx] WARNING: slideIndex {out_of_range_indices} out of range "
            f"(total={total_slides}), skipping",
            file=sys.stderr,
        )

    for slide_index, slide in enumerate(prs.slides):
        slide_changed = False
        slide_edit = slide_edit_map.get(slide_index) or {}
        replacements = slide_edit.get("replaceText") or []

        if effective_hex or palette:
            if recolor_slide_background(slide, effective_hex, palette):
                slide_changed = True
                changed_fills += 1

        add_bullets_list = slide_edit.get("addBullets") or []
        copy_shape_block_action = slide_edit.get("copyShapeBlock")
        repeat_copy_block_action = slide_edit.get("repeatCopyShapeBlock")
        convert_to_cards_action = slide_edit.get("convertToCards")
        copy_layout_ref_action = slide_edit.get("copySlideLayoutFromReference")
        # 挿入済み afterText を追跡（同一テキストが見出し枠・本文枠の両方にある場合の二重挿入防止）
        inserted_after_texts: set[str] = set()
        # replaceText/appendToRun で変更されたシェイプ名を追跡（overflow auto-fit のため）
        text_changed_shapes: set[str] = set()
        for shape in iter_shapes(slide.shapes):
            if effective_hex or palette:
                if recolor_fill(shape, effective_hex, palette):
                    slide_changed = True
                    changed_fills += 1
                if recolor_line(shape, effective_hex, palette):
                    slide_changed = True
                    changed_lines += 1
                if recolor_table(shape, effective_hex, palette):
                    slide_changed = True
                    changed_fills += 1
                if not preserve_text_colors:
                    n = recolor_text_runs(shape, effective_hex, palette)
                    if n > 0:
                        slide_changed = True
                        changed_texts += n
            if font_face and apply_font_face(shape, font_face):
                slide_changed = True
            if replacements and replace_text(shape, replacements):
                slide_changed = True
                text_changed_shapes.add(shape.name)
            if add_bullets_list:
                pending = [
                    item for item in add_bullets_list
                    if (item.get("afterText") or "").strip() not in inserted_after_texts
                ]
                if pending and add_bullets_to_shape(shape, pending, inserted_after_texts):
                    slide_changed = True

        # テキスト増量後の group-aware reflow（フォント縮小・高さ拡張・隣接グループ押し下げ）
        if text_changed_shapes:
            for shape_name in text_changed_shapes:
                warning = reflow_group_after_text_change(
                    slide, shape_name, int(prs.slide_height)
                )
                if warning:
                    layout_warnings.append(f"slide{slide_index + 1}: {warning}")

        # copyShapeBlock: shape ペアコピー（shape ループ外で実行・slide_height を渡してオーバーフロー検知）
        if copy_shape_block_action and copy_shape_block(slide, copy_shape_block_action, int(prs.slide_height)):
            slide_changed = True

        # repeatCopyShapeBlock: 「項目数をNにする」専用
        if repeat_copy_block_action:
            rc_res = repeat_copy_shape_block(slide, repeat_copy_block_action, int(prs.slide_height))
            item_count_slide_results[slide_index] = rc_res
            if rc_res["added"] > 0:
                slide_changed = True
            if rc_res.get("layout_mode") == "tight":
                layout_warnings.append(
                    f"P{slide_index + 1}: 項目数追加でスライド下端マージンを確保できませんでした"
                )

        # copySlideLayoutFromReference: 参照スライドのレイアウトをコピーしてターゲットテキストを流し込む
        if copy_layout_ref_action:
            ref_idx = int(copy_layout_ref_action.get("referenceSlideIndex", -1))
            preserve_text = bool(copy_layout_ref_action.get("preserveTargetText", True))
            if copy_slide_layout_from_reference(prs, slide_index, ref_idx, preserve_text):
                slide_changed = True

        # 画像挿入（特定スライド指定 + slideIndex:-1 の全スライド共通）
        if convert_to_cards_action and convert_slide_to_cards(
            slide,
            convert_to_cards_action,
            int(prs.slide_width),
            int(prs.slide_height),
            effective_hex,
            palette,
        ):
            slide_changed = True

        for img_item in image_insert_map.get(slide_index, []) + all_slides_inserts:
            image_path = img_item.get("imagePath")
            if image_path and Path(image_path).exists():
                try:
                    width_pct = float(img_item.get("widthPct", 15))
                except (TypeError, ValueError):
                    width_pct = 15.0
                if insert_image(
                    slide,
                    image_path,
                    img_item.get("position", "top-right"),
                    width_pct,
                    prs.slide_width,
                    prs.slide_height,
                    near_text=img_item.get("nearText", ""),
                    anchor_side=img_item.get("anchorSide", "right"),
                ):
                    slide_changed = True
                    inserted_images += 1
            else:
                print(f"[edit_pptx] image not found, skipping: {image_path}", file=sys.stderr)

        if slide_changed:
            changed_slides.add(slide_index)

    chars_after = count_all_run_chars(prs)
    prs.save(str(output_path))
    if effective_hex or palette:
        update_theme_colors(output_path, effective_hex, palette)

    # 変更スライドの overflow 候補を検出（決定的チェック）
    overflow_report: list[str] = []
    for si in sorted(changed_slides):
        candidates = detect_text_overflow_candidates(prs.slides[si], int(prs.slide_height))
        if candidates:
            overflow_report.append(f"slide{si}:{','.join(candidates)}")
            print(
                f"[edit_pptx] overflow candidates slide{si}: {candidates}",
                file=sys.stderr,
            )

    result: dict[str, Any] = {
      "changedSlides": len(changed_slides),
      "changedSlideIndices": sorted(changed_slides),
      "totalSlides": len(prs.slides),
      "insertedImages": inserted_images,
      "charsBefore": chars_before,
      "charsAfter": chars_after,
      **({"paletteKey": palette_key} if palette_key else {}),
      **({"changedFills": changed_fills, "changedLines": changed_lines, "changedTexts": changed_texts}
         if (effective_hex or palette) else {}),
    }
    if out_of_range_indices:
        result["outOfRangeSlides"] = out_of_range_indices
    if overflow_report:
        result["overflowCandidates"] = overflow_report
    if layout_warnings:
        result["layoutWarnings"] = layout_warnings
    if item_count_slide_results:
        # slideIndex (int key) を str key に変換して JSON 出力
        result["itemCountResults"] = {str(k): v for k, v in item_count_slide_results.items()}
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
